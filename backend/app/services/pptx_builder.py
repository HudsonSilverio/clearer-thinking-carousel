"""
Build an editable .pptx from carousel data.
Google Slides can import the file via File → Import slides.
"""

import io
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

EMU_PER_PX = 9525          # 914 400 EMU/in ÷ 96 px/in
SLIDE_EMU   = 1080 * EMU_PER_PX

COVER_BG_HEX = "#F5F5F8"
TOP_IMG_H    = round(1080 * 0.65)   # 702 px
BOT_H        = 1080 - TOP_IMG_H     # 378 px
PAD          = 80                   # cover horizontal padding
CONTENT_PAD  = 55
CONTENT_TOP  = 70


# ─── Helpers ─────────────────────────────────────────────────────────────────

def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _px(n: int) -> int:
    return n * EMU_PER_PX


def _set_bg(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _solid_png(hex_color: str) -> bytes:
    """1×1 PNG of a solid colour — used as a stretched picture for rects."""
    h = hex_color.lstrip("#")
    color = (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    buf = io.BytesIO()
    Image.new("RGB", (2, 2), color).save(buf, "PNG")
    return buf.getvalue()


def _add_picture(slide, img_bytes: bytes,
                 left_px: int, top_px: int, w_px: int, h_px: int,
                 send_to_back: bool = False):
    stream = io.BytesIO(img_bytes)
    pic = slide.shapes.add_picture(
        stream,
        Emu(_px(left_px)), Emu(_px(top_px)),
        Emu(_px(w_px)),    Emu(_px(h_px)),
    )
    if send_to_back:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
    return pic


def _add_rect(slide, hex_color: str, left_px: int, top_px: int, w_px: int, h_px: int):
    png = _solid_png(hex_color)
    _add_picture(slide, png, left_px, top_px, w_px, h_px)


def _add_textbox(slide, text: str,
                 left_px: int, top_px: int, w_px: int, h_px: int,
                 font_name: str, size_px: int, bold: bool,
                 hex_color: str, align: str = "left") -> None:
    txBox = slide.shapes.add_textbox(
        Emu(_px(left_px)), Emu(_px(top_px)),
        Emu(_px(w_px)),    Emu(_px(h_px)),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font_name
    f.size = Pt(round(size_px * 0.75, 1))  # px → pt  (96 DPI: 1px = 0.75pt)
    f.bold = bold
    f.color.rgb = _rgb(hex_color)


# ─── Slide builders ───────────────────────────────────────────────────────────

def _cover(prs: Presentation, title: str, cover_bytes: Optional[bytes],
           colors: dict, typography: dict) -> None:
    sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    _set_bg(sl, COVER_BG_HEX)

    # Top area: cover image or solid fill
    if cover_bytes:
        _add_picture(sl, cover_bytes, 0, 0, 1080, TOP_IMG_H, send_to_back=True)
    else:
        _add_rect(sl, colors["slide_bg"], 0, 0, 1080, TOP_IMG_H)

    # Bottom panel (explicit rect so colour is visible even if bg bleeds)
    _add_rect(sl, COVER_BG_HEX, 0, TOP_IMG_H, 1080, BOT_H)

    # Title
    _add_textbox(
        sl, title,
        PAD, TOP_IMG_H + 30, 1080 - PAD * 2, BOT_H - 50,
        typography["headline_font"], typography["headline_size"],
        False, colors["headline"], "center",
    )

    # Accent underline
    _add_rect(sl, colors["progress_bar"], 520 - 20, 1080 - 32, 40, 4)


def _content(prs: Presentation, headline: str, body: str,
             index: int, total: int, colors: dict, typography: dict,
             slide_img_bytes: Optional[bytes] = None) -> None:
    progress_w = max(8, round(1080 * (index + 1) / total))
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, colors["slide_bg"])

    ta = typography.get("text_align", "left")
    hs = typography["headline_size"]
    bs = typography["body_size"]

    # Headline
    _add_textbox(
        sl, headline,
        CONTENT_PAD, CONTENT_TOP, 1080 - CONTENT_PAD * 2, round(hs * 1.25 * 3),
        typography["headline_font"], hs, True, colors["headline"], ta,
    )

    # Body
    body_top = CONTENT_TOP + round(hs * 1.25) + 28
    _add_textbox(
        sl, body,
        CONTENT_PAD, body_top, 1080 - CONTENT_PAD * 2, round(bs * 1.45 * 6),
        typography["body_font"], bs, False, colors["body"], ta,
    )

    # Optional slide image
    if slide_img_bytes:
        img_top = body_top + round(bs * 1.45 * 4) + 20
        img_top = min(img_top, 760)
        _add_picture(sl, slide_img_bytes, CONTENT_PAD, img_top, 1080 - CONTENT_PAD * 2, 240)

    # Progress bar
    _add_rect(sl, colors["progress_bar"], 0, 1072, progress_w, 8)


def _cta(prs: Presentation, cta_bytes: Optional[bytes], colors: dict) -> None:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, "#1A1A2E")
    if cta_bytes:
        _add_picture(sl, cta_bytes, 0, 0, 1080, 1080, send_to_back=True)
    _add_rect(sl, colors["progress_bar"], 0, 1072, 1080, 8)


# ─── Public API ───────────────────────────────────────────────────────────────

def build_pptx(
    title: str,
    takeaways: list[dict],
    colors: dict,
    typography: dict,
    out_path: Path,
    cover_img_bytes: Optional[bytes] = None,
    cta_img_bytes: Optional[bytes] = None,
    slide_img_bytes_list: Optional[list[Optional[bytes]]] = None,
) -> Path:
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_EMU)
    prs.slide_height = Emu(SLIDE_EMU)

    total = len(takeaways) + 2

    _cover(prs, title, cover_img_bytes, colors, typography)

    for i, takeaway in enumerate(takeaways):
        sib = slide_img_bytes_list[i] if (slide_img_bytes_list and i < len(slide_img_bytes_list)) else None
        _content(prs, takeaway["headline"], takeaway.get("body", ""),
                 i + 1, total, colors, typography, sib)

    _cta(prs, cta_img_bytes, colors)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
